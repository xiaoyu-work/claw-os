"""doc — Universal document reader for Claw OS.

Throw any file at it, get structured text back.
"""

import csv
import errno
import io
import json
import os
import stat
import sys

from claw_os_sdk import ai
from cos_runtime import memory, policy


_SUMMARIZE_SYSTEM = (
    "Summarize the document into exactly 5 short bullet lines."
)
_EXPLAIN_SYSTEM = (
    "You are a senior engineer. Explain the supplied content to a curious "
    "user. Keep it under 200 words. Use plain prose, no markdown headings."
)
_REWRITE_SYSTEM = (
    "Rewrite the supplied text following the user's instruction. Return "
    "ONLY the rewritten text — no preamble, no markdown fence, no "
    "commentary. Preserve the language of the input."
)


_MAX_INPUT_CHARS = 100_000


def _read_stdin_or_file(args):
    """Pull text from `--file PATH` or stdin. Used by AI verbs.

    Returns a 3-tuple ``(text, source, instruction_or_err)`` where:
      - On success: ``text`` is a non-empty string, ``source`` describes
        where it came from, ``instruction_or_err`` is the optional
        ``--instruction`` string (or ``None``).
      - On failure: ``text`` is ``None``, ``source`` is ``None``,
        ``instruction_or_err`` is an error ``dict``.
    """
    file_path = None
    instruction = None
    rest = []
    i = 0
    while i < len(args):
        if args[i] == "--file" and i + 1 < len(args):
            file_path = args[i + 1]
            i += 2
        elif args[i] == "--instruction" and i + 1 < len(args):
            instruction = args[i + 1]
            i += 2
        else:
            rest.append(args[i])
            i += 1

    if file_path:
        read_result = cmd_read([file_path])
        if isinstance(read_result, dict) and "error" in read_result:
            return None, None, read_result
        text = read_result.get("content", "") if isinstance(read_result, dict) else ""
        return text, file_path, instruction
    # Fall back to stdin (used by Edit / Term piping a buffer).
    if not sys.stdin.isatty():
        return sys.stdin.read(), "<stdin>", instruction
    if rest:
        return " ".join(rest), None, instruction
    return None, None, {"error": "no input — supply --file PATH or pipe text on stdin"}


def _ai_call(*, text, source, system, max_units):
    """Shared helper for summarize/explain/rewrite — gate → JSON."""
    if not text or not text.strip():
        return {"error": "document produced no extractable text", "source": source}
    if len(text) > _MAX_INPUT_CHARS:
        text = text[:_MAX_INPUT_CHARS]

    policy.require("ai.chat.untrusted", wild=True)

    try:
        response = ai.chat(
            prompt=text,
            origin="external-content",
            system=system,
            max_units=max_units,
        )
    except ai.AiBudgetExceeded as exc:
        return {"error": "AI budget exceeded for this app", "detail": exc.payload}
    except ai.AiSafetyViolation as exc:
        return {"error": "safety violation", "detail": exc.payload}
    except ai.AiDenied as exc:
        return {"error": "AI call denied", "detail": exc.payload}
    except ai.AiUnavailable as exc:
        return {"error": f"AI unavailable: {exc}"}
    except ai.AiError as exc:
        return {"error": str(exc)}

    return {
        "text": response.text,
        "source": source,
        "source_chars": len(text),
        "model": response.model,
        "provider": response.provider,
        "usage": {
            "input_tokens": response.usage.input_tokens,
            "output_tokens": response.usage.output_tokens,
            "units": response.usage.units,
        },
        "budget": {
            "period": response.budget.period,
            "units_used": response.budget.units_used,
            "units_cap": response.budget.units_cap,
        },
        "review": {
            "safety": response.review.safety,
            "prompt_redacted": response.review.prompt_redacted,
        },
    }


class _OpenedDocument:
    def __init__(self, fd, resolved_path, metadata):
        self._fd = fd
        self.resolved_path = resolved_path
        self.metadata = metadata

    def __enter__(self):
        return self

    def __exit__(self, _exc_type, _exc, _traceback):
        self.close()

    def close(self):
        if self._fd is not None:
            os.close(self._fd)
            self._fd = None

    def open_binary(self):
        os.lseek(self._fd, 0, os.SEEK_SET)
        return os.fdopen(os.dup(self._fd), "rb")

    def open_text(self, *, newline=None):
        os.lseek(self._fd, 0, os.SEEK_SET)
        return os.fdopen(
            os.dup(self._fd),
            "r",
            encoding="utf-8",
            newline=newline,
        )


def _resolved_descriptor_path(fd, metadata):
    descriptor_path = f"/proc/self/fd/{fd}"
    try:
        os.readlink(descriptor_path)
    except OSError as exc:
        raise OSError(
            errno.ENOTSUP,
            "descriptor-backed document access requires Linux procfs",
        ) from exc
    if metadata.st_nlink == 0:
        raise OSError(
            errno.ESTALE,
            "document was unlinked before authorization",
        )
    return os.path.realpath(descriptor_path)


def _open_document(path, verb):
    if not hasattr(os, "O_NOFOLLOW"):
        raise OSError(
            errno.ENOTSUP,
            "no-follow document access is unavailable on this platform",
        )
    flags = (
        os.O_RDONLY
        | os.O_NOFOLLOW
        | getattr(os, "O_CLOEXEC", 0)
        | getattr(os, "O_NONBLOCK", 0)
    )
    fd = os.open(path, flags)
    keep_open = False
    try:
        metadata = os.fstat(fd)
        if not stat.S_ISREG(metadata.st_mode):
            raise OSError(errno.EINVAL, "document is not a regular file", path)
        resolved_path = _resolved_descriptor_path(fd, metadata)
        policy.require(verb, path=resolved_path)
        document = _OpenedDocument(fd, resolved_path, metadata)
        keep_open = True
        return document
    finally:
        if not keep_open:
            os.close(fd)


def _open_error(path, exc):
    if exc.errno in (errno.ENOENT, errno.ENOTDIR):
        return {"error": f"file not found: {path}"}
    if exc.errno == errno.ELOOP:
        return {"error": f"refusing to follow document symlink: {path}"}
    if exc.errno == errno.EINVAL:
        return {"error": f"not a regular file: {path}"}
    return {"error": f"failed to open document safely: {path}: {exc}"}


def _open_output(path, *, newline=None):
    if not hasattr(os, "O_NOFOLLOW"):
        raise OSError(
            errno.ENOTSUP,
            "no-follow document output is unavailable on this platform",
        )
    parent = os.path.dirname(path) or "."
    leaf = os.path.basename(path)
    if leaf in ("", ".", ".."):
        raise OSError(errno.EINVAL, "invalid output filename", path)

    parent_flags = (
        os.O_RDONLY
        | getattr(os, "O_DIRECTORY", 0)
        | getattr(os, "O_CLOEXEC", 0)
        | getattr(os, "O_NONBLOCK", 0)
    )
    parent_fd = os.open(parent, parent_flags)
    try:
        parent_metadata = os.fstat(parent_fd)
        if not stat.S_ISDIR(parent_metadata.st_mode):
            raise OSError(errno.ENOTDIR, "output parent is not a directory", parent)
        resolved_parent = _resolved_descriptor_path(parent_fd, parent_metadata)
        resolved_output = os.path.join(resolved_parent, leaf)
        policy.require("fs.write", path=resolved_output)
        output_flags = (
            os.O_WRONLY
            | os.O_CREAT
            | os.O_TRUNC
            | os.O_NOFOLLOW
            | getattr(os, "O_CLOEXEC", 0)
            | getattr(os, "O_NONBLOCK", 0)
        )
        output_fd = os.open(leaf, output_flags, 0o666, dir_fd=parent_fd)
    finally:
        os.close(parent_fd)

    keep_open = False
    try:
        metadata = os.fstat(output_fd)
        if not stat.S_ISREG(metadata.st_mode):
            raise OSError(errno.EINVAL, "output is not a regular file", path)
        output = os.fdopen(
            output_fd,
            "w",
            encoding="utf-8",
            newline=newline,
        )
        keep_open = True
        return output
    finally:
        if not keep_open:
            os.close(output_fd)


def _output_error(path, exc):
    if exc.errno == errno.ELOOP:
        return {"error": f"refusing to follow output symlink: {path}"}
    return {"error": f"failed to open output safely: {path}: {exc}"}


def _read_txt(document):
    with document.open_text() as f:
        text = f.read()
    return text


def _read_json(document):
    with document.open_text() as f:
        data = json.load(f)
    return json.dumps(data, indent=2, ensure_ascii=False)


def _read_csv(document):
    with document.open_text(newline="") as f:
        reader = csv.DictReader(f)
        rows = list(reader)
    return json.dumps(rows, indent=2, ensure_ascii=False)


def _read_yaml(document):
    try:
        import yaml
    except ImportError:
        # yaml is in stdlib via pyyaml on most systems; fall back to raw text
        return _read_txt(document)
    with document.open_text() as f:
        data = yaml.safe_load(f)
    return json.dumps(data, indent=2, ensure_ascii=False)


def _read_pdf(document):
    try:
        import fitz  # pymupdf
    except ImportError:
        return None, {"error": "pymupdf is not installed", "hint": "cos pkg need python3-pymupdf"}
    with document.open_binary() as source:
        doc = fitz.open(stream=source.read(), filetype="pdf")
    try:
        pages = []
        for page in doc:
            pages.append(page.get_text())
    finally:
        doc.close()
    return "\n".join(pages), None


def _read_docx(document):
    try:
        import docx
    except ImportError:
        return None, {"error": "python-docx is not installed", "hint": "cos pkg need python3-docx"}
    with document.open_binary() as source:
        doc = docx.Document(source)
        paragraphs = [p.text for p in doc.paragraphs]
    return "\n".join(paragraphs), None


def _read_xlsx(document):
    try:
        import openpyxl
    except ImportError:
        return None, {"error": "openpyxl is not installed", "hint": "cos pkg need python3-openpyxl"}
    with document.open_binary() as source:
        wb = openpyxl.load_workbook(source, read_only=True, data_only=True)
        try:
            sheets = {}
            for name in wb.sheetnames:
                ws = wb[name]
                rows = []
                for row in ws.iter_rows(values_only=True):
                    rows.append([str(c) if c is not None else "" for c in row])
                sheets[name] = rows
        finally:
            wb.close()
    return json.dumps(sheets, indent=2, ensure_ascii=False), None


def _read_pptx(document):
    try:
        from pptx import Presentation
    except ImportError:
        return None, {"error": "python-pptx is not installed", "hint": "pip install python-pptx"}
    with document.open_binary() as source:
        prs = Presentation(source)
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            texts.append(text)
            notes = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
            slide_text = f"--- Slide {i} ---\n" + "\n".join(texts)
            if notes:
                slide_text += f"\n\n[Notes] {notes}"
            slides.append(slide_text)
    return "\n\n".join(slides), None


def _ext(path):
    return os.path.splitext(path)[1].lower()


def _line_count(text):
    if not text:
        return 0
    return text.count("\n") + (0 if text.endswith("\n") else 1)


# ---------------------------------------------------------------------------
# Commands
# ---------------------------------------------------------------------------

def cmd_read(args):
    if not args:
        return {"error": "usage: cos doc read <path>"}
    path = args[0]
    try:
        document = _open_document(path, "fs.read")
    except OSError as exc:
        return _open_error(path, exc)

    with document:
        ext = _ext(path)
        fmt = ext.lstrip(".") or "txt"

        # Formats that need external libs — may return an error dict
        if ext == ".pdf":
            content, err = _read_pdf(document)
            if err:
                return err
        elif ext == ".docx":
            content, err = _read_docx(document)
            if err:
                return err
        elif ext == ".xlsx":
            content, err = _read_xlsx(document)
            if err:
                return err
        elif ext == ".pptx":
            content, err = _read_pptx(document)
            if err:
                return err
        elif ext in (".yaml", ".yml"):
            content = _read_yaml(document)
        elif ext == ".json":
            try:
                content = _read_json(document)
            except (json.JSONDecodeError, UnicodeDecodeError) as e:
                return {"error": f"failed to parse JSON: {e}"}
        elif ext == ".csv":
            try:
                content = _read_csv(document)
            except (UnicodeDecodeError, csv.Error) as e:
                return {"error": f"failed to parse CSV: {e}"}
        elif ext in (".txt", ".md"):
            try:
                content = _read_txt(document)
            except UnicodeDecodeError:
                return {"error": f"unsupported format: {ext} (binary content)"}
        else:
            # Unknown extension — try reading as UTF-8 text
            try:
                content = _read_txt(document)
            except UnicodeDecodeError:
                return {"error": f"unsupported format: {ext}"}

    return {
        "path": path,
        "format": fmt,
        "content": content,
        "lines": _line_count(content),
    }


def cmd_info(args):
    if not args:
        return {"error": "usage: cos doc info <path>"}
    path = args[0]
    try:
        document = _open_document(path, "fs.meta")
    except OSError as exc:
        return _open_error(path, exc)

    with document:
        ext = _ext(path)
        fmt = ext.lstrip(".") or "txt"
        size = document.metadata.st_size

        # Determine readability: can we handle this format?
        readable = True
        if ext == ".pdf":
            try:
                import fitz  # noqa: F401
            except ImportError:
                readable = False
        elif ext == ".docx":
            try:
                import docx  # noqa: F401
            except ImportError:
                readable = False
        elif ext == ".xlsx":
            try:
                import openpyxl  # noqa: F401
            except ImportError:
                readable = False
        elif ext == ".pptx":
            try:
                from pptx import Presentation  # noqa: F401
            except ImportError:
                readable = False
        else:
            # For unknown extensions, probe if it looks like text.
            if ext not in (".txt", ".md", ".json", ".csv", ".yaml", ".yml"):
                try:
                    with document.open_text() as source:
                        source.read(512)
                except (UnicodeDecodeError, OSError):
                    readable = False

    return {
        "path": path,
        "format": fmt,
        "size": size,
        "readable": readable,
    }


def cmd_summarize(args):
    """Read a document and pipe its text through the AI gate."""
    text, source, extra = _read_stdin_or_file(args)
    if text is None:
        return extra if isinstance(extra, dict) else {"error": "no input"}

    result = _ai_call(
        text=text,
        source=source,
        system=_SUMMARIZE_SYSTEM,
        max_units=6000,
    )
    if "error" in result:
        return result
    out = dict(result)
    # Preserve legacy field name "summary" for backwards compatibility.
    out["summary"] = out.pop("text")
    _remember_doc_summary(source, out.get("summary", ""))
    return out


def _remember_doc_summary(source, summary):
    try:
        if not summary or not source:
            return
        first = summary.strip().splitlines()
        head = first[0] if first else ""
        if len(head) > 200:
            head = head[:197] + "..."
        memory.remember(
            source="doc",
            text=f"Summarised document {source}: {head}",
            kind="note",
            entity_id=source,
            tags=["doc", "summary"],
            link=f"cos doc summarize --file {source}" if source != "<stdin>" else None,
        )
    except memory.MemoryError:
        pass


def cmd_explain(args):
    """Explain the supplied content (via stdin or --file) in plain prose."""
    text, source, extra = _read_stdin_or_file(args)
    if text is None:
        return extra if isinstance(extra, dict) else {"error": "no input"}
    return _ai_call(
        text=text,
        source=source,
        system=_EXPLAIN_SYSTEM,
        max_units=4000,
    )


def cmd_rewrite(args):
    """Rewrite the supplied content following an `--instruction`."""
    text, source, instruction = _read_stdin_or_file(args)
    if text is None:
        return instruction if isinstance(instruction, dict) else {"error": "no input"}
    if not instruction:
        instruction = "Improve clarity, fix grammar, keep the original meaning."
    system = _REWRITE_SYSTEM + "\n\nInstruction: " + instruction
    return _ai_call(
        text=text,
        source=source,
        system=system,
        max_units=8000,
    )


def cmd_convert(args):
    if not args:
        return {"error": "usage: cos doc convert <path> --to <format>"}

    path = args[0]
    target_fmt = None
    for i, a in enumerate(args):
        if a == "--to" and i + 1 < len(args):
            target_fmt = args[i + 1].lstrip(".")
            break
    if not target_fmt:
        return {"error": "usage: cos doc convert <path> --to <format>"}

    ext = _ext(path)
    if (ext, target_fmt) not in ((".json", "csv"), (".csv", "json")):
        return {"error": "unsupported conversion"}
    base = os.path.splitext(path)[0]
    output_path = f"{base}.{target_fmt}"

    try:
        document = _open_document(path, "fs.read")
    except OSError as exc:
        return _open_error(path, exc)

    # JSON -> CSV
    if ext == ".json" and target_fmt == "csv":
        with document:
            try:
                with document.open_text() as source:
                    data = json.load(source)
            except (json.JSONDecodeError, UnicodeDecodeError) as e:
                return {"error": f"failed to parse JSON: {e}"}
        if not isinstance(data, list) or not data or not isinstance(data[0], dict):
            return {"error": "JSON must be an array of objects for CSV conversion"}
        fieldnames = list(data[0].keys())
        buf = io.StringIO()
        writer = csv.DictWriter(buf, fieldnames=fieldnames)
        writer.writeheader()
        for row in data:
            writer.writerow(row)
        try:
            with _open_output(output_path, newline="") as output:
                output.write(buf.getvalue())
        except OSError as exc:
            return _output_error(output_path, exc)
        return {"input": path, "output": output_path, "format": "csv"}

    # CSV -> JSON
    with document:
        try:
            with document.open_text(newline="") as source:
                reader = csv.DictReader(source)
                rows = list(reader)
        except (UnicodeDecodeError, csv.Error) as e:
            return {"error": f"failed to parse CSV: {e}"}
    try:
        with _open_output(output_path) as output:
            json.dump(rows, output, indent=2, ensure_ascii=False)
    except OSError as exc:
        return _output_error(output_path, exc)
    return {"input": path, "output": output_path, "format": "json"}


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

def _schema():
    return {
        "read": {
            "description": "Read a document and return structured text (supports txt, md, json, csv, yaml, pdf, docx, xlsx, pptx)",
            "parameters": [
                {"name": "path", "type": "string", "required": True, "description": "Path to the document file", "kind": "positional"},
            ],
            "example": "cos app doc read /workspace/report.pdf",
        },
        "info": {
            "description": "Get document metadata (format, size, readability)",
            "parameters": [
                {"name": "path", "type": "string", "required": True, "description": "Path to the document file", "kind": "positional"},
            ],
            "example": "cos app doc info /workspace/report.pdf",
        },
        "convert": {
            "description": "Convert between document formats (JSON to CSV, CSV to JSON)",
            "parameters": [
                {"name": "path", "type": "string", "required": True, "description": "Path to the source file", "kind": "positional"},
                {"name": "--to", "type": "string", "required": True, "description": "Target format (e.g. csv, json)", "kind": "flag"},
            ],
            "example": "cos app doc convert /workspace/data.json --to csv",
        },
        "summarize": {
            "description": "Summarize a document into 5 short bullet lines via the AI gate. Reads from --file or stdin.",
            "parameters": [
                {"name": "--file", "type": "string", "required": False, "description": "Path to the document to summarise (or omit and pipe via stdin)", "kind": "flag"},
            ],
            "example": "cos app doc summarize --file /workspace/report.pdf",
        },
        "explain": {
            "description": "Explain the supplied content (via --file or stdin) in plain prose under 200 words.",
            "parameters": [
                {"name": "--file", "type": "string", "required": False, "description": "Path to the document to explain (or pipe via stdin)", "kind": "flag"},
            ],
            "example": "echo 'fn foo()...' | cos app doc explain",
        },
        "rewrite": {
            "description": "Rewrite the supplied content per --instruction. Returns rewritten text only.",
            "parameters": [
                {"name": "--file", "type": "string", "required": False, "description": "Path to read (or pipe via stdin)", "kind": "flag"},
                {"name": "--instruction", "type": "string", "required": False, "description": "Rewrite instruction (e.g. 'make it more formal')", "kind": "flag"},
            ],
            "example": "cos app doc rewrite --file note.md --instruction 'translate to English'",
        },
    }


def run(command, args):
    """Called by cos router."""
    if command == "__schema__":
        return _schema()
    commands = {
        "read": cmd_read,
        "info": cmd_info,
        "convert": cmd_convert,
        "summarize": cmd_summarize,
        "explain": cmd_explain,
        "rewrite": cmd_rewrite,
    }
    handler = commands.get(command)
    if not handler:
        return {"error": f"unknown command: {command}"}
    try:
        return handler(args)
    except policy.PermissionDenied as denied:
        return {"error": str(denied), "denial": denied.denial}
    except policy.PolicyUnavailable as exc:
        return {"error": f"capability check failed: {exc}"}
