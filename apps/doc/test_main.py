import os
import sys
import types
from pathlib import Path

import pytest

from test_support import load_local_module


pytestmark = pytest.mark.skipif(
    sys.platform != "linux" or not hasattr(os, "O_NOFOLLOW"),
    reason="the document app's descriptor-safe contract targets Claw OS Linux",
)


@pytest.fixture
def doc_app():
    return load_local_module(Path(__file__).with_name("main.py"), "doc_app_main")


def _swap_during_authorization(
    monkeypatch,
    doc_app,
    requested,
    replacement,
    *,
    verb="fs.read",
):
    authorized_path = os.path.realpath(requested)

    def require(actual_verb, *, path=None, **_scope):
        assert actual_verb == verb
        assert path == authorized_path
        requested.unlink()
        requested.symlink_to(replacement)

    monkeypatch.setattr(doc_app.policy, "require", require)


def test_read_keeps_opened_file_when_path_becomes_symlink(
    tmp_path,
    monkeypatch,
    doc_app,
):
    requested = tmp_path / "note.txt"
    replacement = tmp_path / "secret.txt"
    requested.write_text("authorized content", encoding="utf-8")
    replacement.write_text("secret content", encoding="utf-8")
    _swap_during_authorization(
        monkeypatch,
        doc_app,
        requested,
        replacement,
    )

    result = doc_app.cmd_read([str(requested)])

    assert result["content"] == "authorized content"
    assert requested.read_text(encoding="utf-8") == "secret content"


def test_read_rejects_final_symlink_before_policy(
    tmp_path,
    monkeypatch,
    doc_app,
):
    secret = tmp_path / "secret.txt"
    link = tmp_path / "link.txt"
    secret.write_text("secret content", encoding="utf-8")
    link.symlink_to(secret)

    def unexpected_authorization(*_args, **_kwargs):
        pytest.fail("a final symlink must be rejected before authorization")

    monkeypatch.setattr(doc_app.policy, "require", unexpected_authorization)

    result = doc_app.cmd_read([str(link)])

    assert "symlink" in result["error"]
    assert "secret content" not in result["error"]


def test_info_uses_fstat_from_authorized_descriptor(
    tmp_path,
    monkeypatch,
    doc_app,
):
    requested = tmp_path / "report.unknown"
    replacement = tmp_path / "secret.unknown"
    requested.write_text("short", encoding="utf-8")
    replacement.write_text("secret content that is longer", encoding="utf-8")
    _swap_during_authorization(
        monkeypatch,
        doc_app,
        requested,
        replacement,
        verb="fs.meta",
    )

    result = doc_app.cmd_info([str(requested)])

    assert result["size"] == len("short")
    assert result["readable"] is True


@pytest.mark.parametrize(
    ("suffix", "authorized", "secret"),
    [
        (".txt", "authorized text", "secret text"),
        (".md", "authorized markdown", "secret markdown"),
        (".json", '{"value": "authorized json"}', '{"value": "secret json"}'),
        (".csv", "value\nauthorized csv\n", "value\nsecret csv\n"),
        (".yaml", "value: authorized yaml\n", "value: secret yaml\n"),
        (".unknown", "authorized fallback", "secret fallback"),
    ],
)
def test_text_readers_share_authorized_descriptor_contract(
    tmp_path,
    monkeypatch,
    doc_app,
    suffix,
    authorized,
    secret,
):
    requested = tmp_path / f"document{suffix}"
    replacement = tmp_path / f"secret{suffix}"
    requested.write_text(authorized, encoding="utf-8")
    replacement.write_text(secret, encoding="utf-8")
    _swap_during_authorization(
        monkeypatch,
        doc_app,
        requested,
        replacement,
    )

    result = doc_app.cmd_read([str(requested)])

    assert "authorized" in result["content"]
    assert "secret" not in result["content"]


def test_packaged_readers_share_authorized_descriptor_contract(
    tmp_path,
    monkeypatch,
    doc_app,
):
    authorized = b"authorized package bytes"
    secret = b"secret package bytes"

    class FakePdf:
        def __iter__(self):
            return iter([types.SimpleNamespace(get_text=lambda: "authorized pdf")])

        def close(self):
            return None

    def open_pdf(*, stream, filetype):
        assert stream == authorized
        assert filetype == "pdf"
        return FakePdf()

    monkeypatch.setitem(sys.modules, "fitz", types.SimpleNamespace(open=open_pdf))

    def open_docx(source):
        assert source.read() == authorized
        return types.SimpleNamespace(
            paragraphs=[types.SimpleNamespace(text="authorized docx")]
        )

    monkeypatch.setitem(
        sys.modules,
        "docx",
        types.SimpleNamespace(Document=open_docx),
    )

    class FakeWorkbook:
        sheetnames = ["Sheet1"]

        def __getitem__(self, _name):
            return types.SimpleNamespace(
                iter_rows=lambda **_kwargs: iter([("authorized xlsx",)])
            )

        def close(self):
            return None

    def open_xlsx(source, *, read_only, data_only):
        assert source.read() == authorized
        assert read_only is True
        assert data_only is True
        return FakeWorkbook()

    monkeypatch.setitem(
        sys.modules,
        "openpyxl",
        types.SimpleNamespace(load_workbook=open_xlsx),
    )

    def open_pptx(source):
        assert source.read() == authorized
        paragraph = types.SimpleNamespace(text="authorized pptx")
        shape = types.SimpleNamespace(
            has_text_frame=True,
            text_frame=types.SimpleNamespace(paragraphs=[paragraph]),
        )
        slide = types.SimpleNamespace(
            shapes=[shape],
            has_notes_slide=False,
        )
        return types.SimpleNamespace(slides=[slide])

    monkeypatch.setitem(
        sys.modules,
        "pptx",
        types.SimpleNamespace(Presentation=open_pptx),
    )

    for suffix in (".pdf", ".docx", ".xlsx", ".pptx"):
        requested = tmp_path / f"document{suffix}"
        replacement = tmp_path / f"secret{suffix}"
        requested.write_bytes(authorized)
        replacement.write_bytes(secret)
        _swap_during_authorization(
            monkeypatch,
            doc_app,
            requested,
            replacement,
        )

        result = doc_app.cmd_read([str(requested)])

        assert "authorized" in result["content"]
        assert "secret" not in result["content"]


def test_real_packaged_readers_survive_symlink_swap(
    tmp_path,
    monkeypatch,
    doc_app,
):
    fitz = pytest.importorskip("fitz")
    docx = pytest.importorskip("docx")
    openpyxl = pytest.importorskip("openpyxl")
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches

    def write_pdf(path, text):
        document = fitz.open()
        page = document.new_page()
        page.insert_text((72, 72), text)
        document.save(path)
        document.close()

    def write_docx(path, text):
        document = docx.Document()
        document.add_paragraph(text)
        document.save(path)

    def write_xlsx(path, text):
        workbook = openpyxl.Workbook()
        workbook.active["A1"] = text
        workbook.save(path)
        workbook.close()

    def write_pptx(path, text):
        presentation = pptx.Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        box = slide.shapes.add_textbox(
            Inches(1),
            Inches(1),
            Inches(5),
            Inches(1),
        )
        box.text = text
        presentation.save(path)

    writers = {
        ".pdf": write_pdf,
        ".docx": write_docx,
        ".xlsx": write_xlsx,
        ".pptx": write_pptx,
    }
    for suffix, writer in writers.items():
        requested = tmp_path / f"real{suffix}"
        replacement = tmp_path / f"secret-real{suffix}"
        writer(requested, f"authorized {suffix}")
        writer(replacement, f"secret {suffix}")
        _swap_during_authorization(
            monkeypatch,
            doc_app,
            requested,
            replacement,
        )

        result = doc_app.cmd_read([str(requested)])

        assert "authorized" in result["content"]
        assert "secret" not in result["content"]


def test_convert_reads_opened_source_after_symlink_swap(
    tmp_path,
    monkeypatch,
    doc_app,
):
    requested = tmp_path / "records.json"
    replacement = tmp_path / "secret.json"
    requested.write_text('[{"value": "authorized"}]', encoding="utf-8")
    replacement.write_text('[{"value": "secret"}]', encoding="utf-8")
    authorized_path = os.path.realpath(requested)
    swapped = False

    def require(verb, *, path=None, **_scope):
        nonlocal swapped
        if verb == "fs.read":
            assert path == authorized_path
            requested.unlink()
            requested.symlink_to(replacement)
            swapped = True
        else:
            assert verb == "fs.write"

    monkeypatch.setattr(doc_app.policy, "require", require)

    result = doc_app.cmd_convert([str(requested), "--to", "csv"])

    assert swapped is True
    assert result["format"] == "csv"
    output = Path(result["output"]).read_text(encoding="utf-8")
    assert "authorized" in output
    assert "secret" not in output


def test_convert_rejects_output_symlink_after_authorization(
    tmp_path,
    monkeypatch,
    doc_app,
):
    source = tmp_path / "records.json"
    outside = tmp_path / "outside.csv"
    output = tmp_path / "records.csv"
    source.write_text('[{"value": "authorized"}]', encoding="utf-8")
    outside.write_text("do not overwrite", encoding="utf-8")
    output.symlink_to(outside)
    monkeypatch.setattr(doc_app.policy, "require", lambda *_args, **_kwargs: None)

    result = doc_app.cmd_convert([str(source), "--to", "csv"])

    assert "output symlink" in result["error"]
    assert outside.read_text(encoding="utf-8") == "do not overwrite"
